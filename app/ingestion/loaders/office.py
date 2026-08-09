import logfire
import os
from docx import Document
from pptx import Presentation

def parse_office(file_path: str) -> str:
    """
    Parses Office documents (.docx, .pptx) using python-docx and python-pptx.
    These are structured and lightweight, eliminating the need for ML-heavy libraries.
    """
    with logfire.span("📄 Office Document Parsing", filename=file_path):
        try:
            ext = file_path.lower().rsplit(".", 1)[-1]
            full_text = ""

            if ext == "docx":
                doc = Document(file_path)
                # Extract text from paragraphs, ignoring empty ones
                full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
                
            elif ext == "pptx":
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        # Only extract from shapes that actually contain text
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                            text_parts.append(shape.text.strip())
                full_text = "\n".join([part for part in text_parts if part])
                
            else:
                logfire.warning(f"⚠️ Unsupported office extension '{ext}' for {file_path}")
                return ""

            if not full_text.strip():
                logfire.warning(f"⚠️ Parser returned empty text for {file_path}")
            else:
                logfire.info(f"✅ Successfully parsed {len(full_text)} characters")

            return full_text
            
        except Exception as e:
            logfire.error(f"❌ Office Parse Failed: {e}")
            raise e